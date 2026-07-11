"""
Phase 11 — extract a slide deck from an uploaded .pptx.

A .pptx is an Office Open XML package (a zip of XML). `python-pptx` reads the
text for us. We pull, per slide: the title, the body bullet points, any table
text, and the speaker notes — into the SAME {title, bullets[], note} structure
the rest of the app already uses (so Q&A, voice, and Present Mode just work).

Constraints (see chat notes): images/diagrams/charts are opaque to text
extraction; only text + notes are read. Long decks are capped for token safety.
"""

from io import BytesIO

from pptx import Presentation

MAX_UPLOAD_SLIDES = 12      # cap for token/context safety
MAX_BULLETS = 6
_TITLE_PH = {"title", "ctrTitle", "subTitle"}


def _para_text(paragraph) -> str:
    txt = "".join(run.text for run in paragraph.runs).strip()
    return txt or (paragraph.text or "").strip()


def extract_deck(data: bytes) -> dict:
    prs = Presentation(BytesIO(data))
    slides = []

    for slide in prs.slides:
        title_shape = slide.shapes.title
        title = (title_shape.text or "").strip() if title_shape is not None else ""
        # python-pptx makes fresh wrapper objects per access, so compare by id.
        title_id = title_shape.shape_id if title_shape is not None else None
        bullets = []

        for shape in slide.shapes:
            if title_id is not None and shape.shape_id == title_id:
                continue
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        bullets.append(" | ".join(cells))
                continue
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = _para_text(para)
                if t:
                    bullets.append(t)

        note = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            note = (slide.notes_slide.notes_text_frame.text or "").strip()

        # If the slide had no title placeholder, promote the first line.
        if not title:
            if bullets:
                title, bullets = bullets[0], bullets[1:]
            else:
                title = f"Slide {len(slides) + 1}"

        bullets = [b for b in bullets if b][:MAX_BULLETS]
        slides.append({
            "title": title[:140],
            "bullets": bullets or ["(no text on this slide)"],
            "note": note[:300],
        })

    truncated = len(slides) > MAX_UPLOAD_SLIDES
    slides = slides[:MAX_UPLOAD_SLIDES]

    title = "Presentation"
    try:
        if prs.core_properties.title:
            title = prs.core_properties.title.strip()
    except Exception:
        pass
    if title == "Presentation" and slides:
        title = slides[0]["title"]

    return {"title": title, "slides": slides, "truncated": truncated}
